import pptx
import os

def pptx_to_text(pptx_file):
    try:
        prs = pptx.Presentation(pptx_file)
        text_file = os.path.splitext(pptx_file)[0] + ".txt"
        with open(text_file, "w") as f:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        f.write(shape.text)
                        f.write("\n")
        print(f"Successfully converted {pptx_file} to {text_file}")
    except Exception as e:
        print(f"An error occurred: {e}")

if __name__ == "__main__":
    pptx_to_text("Viterbi_Decoder.pptx")
